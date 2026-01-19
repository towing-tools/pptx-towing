#!/usr/bin/env python3
"""
Extract content from PPTX file for analysis.
Outputs slide structure and content to markdown format.
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

def extract_shape_info(shape, indent=0):
    """Extract information from a shape."""
    info = []
    prefix = "  " * indent

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                # Get font info
                font_size = None
                bold = False
                for run in paragraph.runs:
                    if run.font.size:
                        font_size = run.font.size.pt
                    if run.font.bold:
                        bold = True

                font_info = f" (size: {font_size}pt, bold: {bold})" if font_size else ""
                info.append(f"{prefix}- Text: \"{text}\"{font_info}")

    if shape.has_table:
        table = shape.table
        info.append(f"{prefix}- Table ({table.rows.__len__()}x{table.columns.__len__()}):")
        for row_idx, row in enumerate(table.rows):
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip() if cell.text else "")
            info.append(f"{prefix}  Row {row_idx + 1}: {' | '.join(row_data)}")

    if hasattr(shape, 'shapes'):  # Group shape
        info.append(f"{prefix}- Group Shape:")
        for child in shape.shapes:
            info.extend(extract_shape_info(child, indent + 1))

    return info

def analyze_pptx(pptx_path):
    """Analyze PPTX file and extract structure."""
    prs = Presentation(pptx_path)

    output = []
    output.append(f"# PPTX Analysis: {os.path.basename(pptx_path)}")
    output.append(f"")
    output.append(f"## Overview")
    output.append(f"- Total Slides: {len(prs.slides)}")
    output.append(f"- Slide Width: {prs.slide_width.inches:.2f} inches")
    output.append(f"- Slide Height: {prs.slide_height.inches:.2f} inches")
    output.append(f"")

    for slide_idx, slide in enumerate(prs.slides, 1):
        output.append(f"---")
        output.append(f"")
        output.append(f"## Slide {slide_idx}")
        output.append(f"")

        # Get slide layout name if available
        try:
            layout_name = slide.slide_layout.name
            output.append(f"**Layout**: {layout_name}")
        except:
            pass

        output.append(f"**Shapes**: {len(slide.shapes)}")
        output.append(f"")

        for shape_idx, shape in enumerate(slide.shapes, 1):
            shape_type = "Unknown"
            try:
                shape_type = shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type)
            except:
                pass

            # Position info
            pos_info = ""
            try:
                pos_info = f" @ ({shape.left.inches:.2f}\", {shape.top.inches:.2f}\") size ({shape.width.inches:.2f}\"x{shape.height.inches:.2f}\")"
            except:
                pass

            output.append(f"### Shape {shape_idx}: {shape_type}{pos_info}")

            # Fill color
            try:
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    fill_type = str(shape.fill.type)
                    output.append(f"- Fill: {fill_type}")
                    if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color.type is not None:
                        try:
                            rgb = shape.fill.fore_color.rgb
                            output.append(f"- Color: #{rgb}")
                        except:
                            pass
            except:
                pass

            # Extract content
            shape_info = extract_shape_info(shape)
            if shape_info:
                output.extend(shape_info)

            output.append(f"")

    return "\n".join(output)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx_content.py <pptx_file> [output_file]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    result = analyze_pptx(pptx_path)

    if output_path:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(result)
        print(f"Analysis saved to: {output_path}")
    else:
        print(result)
